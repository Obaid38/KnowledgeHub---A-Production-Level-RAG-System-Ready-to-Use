import io
import logging
import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from ai.config import OCR_DPI, OCR_LANGUAGES, PDF_MIN_CHARS_BEFORE_OCR_FALLBACK, TEMP_DIR

logger = logging.getLogger("knowledge_hub.extractor")


@dataclass
class ExtractionResult:
    text: str
    page_count: int = 0
    extraction_method: str = "text_extraction"
    raw_tables: list[list[list[str]]] = field(default_factory=list)


def extract(file_bytes: bytes, filename: str) -> ExtractionResult:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_bytes)
    elif ext == ".xlsx":
        return _extract_xlsx(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes)
    elif ext in {".jpg", ".jpeg", ".png", ".tiff"}:
        return _extract_image(file_bytes)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _extract_pdf(file_bytes: bytes) -> ExtractionResult:
    import fitz  # PyMuPDF

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf", dir=str(TEMP_DIR))
    try:
        tmp.write(file_bytes)
        tmp.flush()
        tmp.close()  # close write handle before fitz/pdfplumber open (required on Windows)

        # Text extraction via PyMuPDF
        doc = fitz.open(str(tmp.name))
        page_count = len(doc)
        pages_text = []
        for page in doc:
            pages_text.append(page.get_text())
        doc.close()
        full_text = "\n".join(pages_text)
        char_count = len(full_text.strip())

        # OCR fallback if insufficient text
        if char_count < PDF_MIN_CHARS_BEFORE_OCR_FALLBACK:
            logger.info(
                "[Extractor] PDF text too short (%d chars), falling back to OCR", char_count
            )
            return _ocr_pdf(Path(tmp.name), page_count)

        # Table extraction via pdfplumber
        raw_tables: list[list[list[str]]] = []
        try:
            import pdfplumber

            with pdfplumber.open(str(tmp.name)) as pdf:
                for page in pdf.pages:
                    for table in page.extract_tables() or []:
                        if table:
                            # Normalise: each cell to str
                            norm = [
                                [str(cell) if cell is not None else "" for cell in row]
                                for row in table
                            ]
                            raw_tables.append(norm)
        except Exception as exc:
            logger.warning("[Extractor] pdfplumber table extraction failed: %s", exc)

        logger.info(
            "[Extractor] PDF '%s': %d pages, %d chars, %d tables",
            "...",
            page_count,
            char_count,
            len(raw_tables),
        )
        return ExtractionResult(
            text=full_text,
            page_count=page_count,
            extraction_method="text_extraction",
            raw_tables=raw_tables,
        )
    finally:
        try:
            os.unlink(tmp.name)
        except OSError:
            pass


def _ocr_pdf(pdf_path: Path, page_count: int) -> ExtractionResult:
    import fitz
    from PIL import Image
    import pytesseract

    doc = fitz.open(str(pdf_path))
    texts = []
    for page in doc:
        mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        text = pytesseract.image_to_string(img, lang=OCR_LANGUAGES)
        texts.append(text)
    doc.close()

    return ExtractionResult(
        text="\n".join(texts),
        page_count=page_count,
        extraction_method="ocr",
        raw_tables=[],
    )


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_docx(file_bytes: bytes) -> ExtractionResult:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn

    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs: list[str] = []
    raw_tables: list[list[list[str]]] = []

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
        if tag == "p":
            text = "".join(t.text for t in block.iter(qn("w:t")) if t.text is not None)
            paragraphs.append(text)
        elif tag == "tbl":
            table_rows: list[list[str]] = []
            for row in block.iter(qn("w:tr")):
                cells = []
                for cell in row.iter(qn("w:tc")):
                    cell_text = "".join(t.text for t in cell.iter(qn("w:t")) if t.text is not None)
                    cells.append(cell_text)
                if cells:
                    table_rows.append(cells)
            if table_rows:
                raw_tables.append(table_rows)

    full_text = "\n".join(paragraphs)
    logger.info("[Extractor] DOCX: %d chars, %d tables", len(full_text), len(raw_tables))
    return ExtractionResult(
        text=full_text,
        page_count=0,
        extraction_method="structured_parse",
        raw_tables=raw_tables,
    )


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def _extract_xlsx(file_bytes: bytes) -> ExtractionResult:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    raw_tables: list[list[list[str]]] = []
    text_lines: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_lines.append(f"[Sheet: {sheet_name}]")
        table_rows: list[list[str]] = []
        for row in ws.iter_rows():
            cells = [str(c.value) if c.value is not None else "" for c in row]
            if any(cells):
                table_rows.append(cells)
                text_lines.append("\t".join(cells))
        if table_rows:
            raw_tables.append(table_rows)

    wb.close()
    full_text = "\n".join(text_lines)
    logger.info("[Extractor] XLSX: %d chars, %d sheets", len(full_text), len(wb.sheetnames) if hasattr(wb, 'sheetnames') else 0)
    return ExtractionResult(
        text=full_text,
        page_count=0,
        extraction_method="structured_parse",
        raw_tables=raw_tables,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _extract_pptx(file_bytes: bytes) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slide_count = len(prs.slides)
    text_lines: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_lines.append(f"[Slide {slide_idx}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        text_lines.append(line)

    full_text = "\n".join(text_lines)
    logger.info("[Extractor] PPTX: %d slides, %d chars", slide_count, len(full_text))
    return ExtractionResult(
        text=full_text,
        page_count=slide_count,
        extraction_method="structured_parse",
        raw_tables=[],
    )


# ---------------------------------------------------------------------------
# Images (OCR)
# ---------------------------------------------------------------------------

def _extract_image(file_bytes: bytes) -> ExtractionResult:
    from PIL import Image
    import pytesseract

    img = Image.open(io.BytesIO(file_bytes))
    text = pytesseract.image_to_string(img, lang=OCR_LANGUAGES, config=f"--dpi {OCR_DPI}")
    logger.info("[Extractor] Image OCR: %d chars", len(text))
    return ExtractionResult(
        text=text,
        page_count=1,
        extraction_method="ocr",
        raw_tables=[],
    )
